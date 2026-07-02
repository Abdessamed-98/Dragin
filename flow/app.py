import os, sys, io, base64, subprocess, threading, uuid, tempfile, shutil, json, zipfile, time
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from PIL import Image, ImageOps, ImageEnhance, ImageFilter
# Register HEIC/HEIF support so Image.open() can read iPhone photos (input-only).
try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
except Exception as _heif_err:
    print(f"[Backend] pillow-heif unavailable, HEIC input disabled: {_heif_err}")
import vtracer

# Defer the heaviest imports (fitz ~90ms, numpy ~84ms) until first actual use — keeps
# backend cold-start fast. This proxy is transparent: every `fitz.`/`np.` call site works
# unchanged; the real module loads on the first attribute access.
class _LazyModule:
    def __init__(self, name):
        self.__dict__['_name'] = name
        self.__dict__['_mod'] = None
    def __getattr__(self, attr):
        mod = self.__dict__['_mod']
        if mod is None:
            import importlib
            mod = importlib.import_module(self.__dict__['_name'])
            self.__dict__['_mod'] = mod
        return getattr(mod, attr)

fitz = _LazyModule('fitz')   # PyMuPDF
np = _LazyModule('numpy')

app = Flask(__name__)
CORS(app)

# ── On-demand tools directory (set by Electron main process) ──────
TOOLS_DIR = os.environ.get('DRAGIN_TOOLS_DIR', '')

print("[Backend] Device: CPU")

# Shared idle timeout for on-demand models (e.g. BEN2) — unload after 5 min idle.
_MODEL_IDLE_TIMEOUT = 300



# ── BEN2 background remover (on-demand, single model, no modes) ──────
# BEN2 (Background Erase Network 2) — higher-quality matting via the `ben2`
# PyTorch package. Point the HuggingFace cache at the downloaded tool dir so
# the weights load offline.
if TOOLS_DIR:
    _ben2_dir = os.path.join(TOOLS_DIR, 'remover')
    if os.path.isdir(_ben2_dir):
        os.environ.setdefault('HF_HOME', os.path.join(_ben2_dir, 'hf'))

_ben2_model = None              # (model, device) tuple once loaded
_ben2_lock = threading.Lock()
_ben2_last_used = 0.0

def _load_ben2():
    """Lazy-load the BEN2 model. Raises ImportError if ben2/torch not installed."""
    global _ben2_model, _ben2_last_used
    with _ben2_lock:
        _ben2_last_used = time.time()
        if _ben2_model is None:
            import torch
            from ben2 import BEN_Base
            print("[BEN2] Loading model ...")
            # Prefer CUDA, then Apple Silicon MPS, else CPU — so Macs aren't stuck on slow CPU.
            if torch.cuda.is_available():
                device = torch.device('cuda')
            elif getattr(torch.backends, 'mps', None) is not None and torch.backends.mps.is_available():
                device = torch.device('mps')
            else:
                device = torch.device('cpu')
            m = BEN_Base.from_pretrained("PramaLLC/BEN2")
            try:
                m.to(device).eval()
            except Exception as e:
                print(f"[BEN2] {device} failed ({e}); falling back to CPU")
                device = torch.device('cpu')
                m.to(device).eval()
            _ben2_model = (m, device)
            print(f"[BEN2] Model ready on {device}.")
        return _ben2_model

def _unload_ben2():
    """Free the cached BEN2 model to reclaim memory."""
    global _ben2_model
    with _ben2_lock:
        if _ben2_model is not None:
            _ben2_model = None
            import gc; gc.collect()
            try:
                import torch
                if torch.cuda.is_available():
                    torch.cuda.empty_cache()
            except Exception:
                pass
            print("[BEN2] Unloaded model — memory freed")

def _ben2_idle_watcher():
    """Background thread: unload BEN2 after idle timeout."""
    global _ben2_model
    while True:
        time.sleep(60)
        with _ben2_lock:
            if _ben2_model is not None and _ben2_last_used > 0:
                if time.time() - _ben2_last_used >= _MODEL_IDLE_TIMEOUT:
                    _ben2_model = None
                    import gc; gc.collect()
                    print("[BEN2] Auto-unloaded after idle")

threading.Thread(target=_ben2_idle_watcher, daemon=True).start()


@app.route('/ben2/process', methods=['POST'])
def ben2_process():
    global _ben2_last_used
    try:
        model, device = _load_ben2()
    except ImportError:
        return jsonify({"error": "BEN2 not installed. Install 'Background Remover (BEN2)' from Store."}), 503

    if 'images' not in request.files:
        return jsonify({"error": "No images"}), 400

    files = request.files.getlist('images')
    results = []
    req_start = time.perf_counter()

    for file in files:
        t0 = time.perf_counter()
        input_image = Image.open(io.BytesIO(file.read())).convert("RGB")
        w, h = input_image.size
        t_load = time.perf_counter() - t0

        t1 = time.perf_counter()
        # refine_foreground is an iterative full-res matte refine — very slow on CPU,
        # so keep it off for the CPU path (the base BEN2 matte is already high quality).
        foreground = model.inference(input_image, refine_foreground=False)
        t_infer = time.perf_counter() - t1

        t2 = time.perf_counter()
        buffered = io.BytesIO()
        foreground.save(buffered, format="PNG")
        img_str = base64.b64encode(buffered.getvalue()).decode('utf-8')
        t_encode = time.perf_counter() - t2

        print(f"[BEN2] {file.filename} ({w}x{h}) device={device} | "
              f"load={t_load:.2f}s infer={t_infer:.2f}s encode={t_encode:.2f}s")

        results.append({"name": file.filename, "data": img_str})

    _ben2_last_used = time.time()  # keep alive across long batches
    print(f"[BEN2] Request done: {len(files)} file(s) in {time.perf_counter()-req_start:.2f}s")
    return jsonify({"results": results})


@app.route('/ben2/model-status', methods=['GET'])
def ben2_model_status():
    """Check whether the BEN2 model is already loaded."""
    with _ben2_lock:
        loaded = _ben2_model is not None
    return jsonify({"loaded": loaded, "model": "BEN2"})


@app.route('/ben2/unload', methods=['POST'])
def ben2_unload():
    """Explicitly free the cached BEN2 model to reclaim RAM."""
    _unload_ben2()
    return jsonify({"ok": True})

@app.route('/compress', methods=['POST'])
def compress():
    if 'image' not in request.files:
        return jsonify({"error": "No image file provided"}), 400

    file = request.files['image']
    quality = int(request.form.get('quality', 70))

    try:
        raw_bytes = file.read()
        original_size = len(raw_bytes)
        input_image = Image.open(io.BytesIO(raw_bytes))
        original_format = (input_image.format or 'PNG').upper()

        output_buffer = io.BytesIO()

        if original_format == 'JPEG' or original_format == 'JPG':
            # JPEG: re-encode at lower quality with optimization
            if input_image.mode in ('RGBA', 'P'):
                input_image = input_image.convert('RGB')
            input_image.save(output_buffer, format='JPEG', quality=quality, optimize=True)
            output_format = 'JPEG'
            mime = 'image/jpeg'
        elif original_format == 'PNG':
            # PNG: map quality (10-95) → color count (16-256) so the slider has real effect
            colors = max(16, min(256, round(quality * 256 / 95)))
            if input_image.mode == 'RGBA':
                quantized = input_image.quantize(colors=colors, method=2, dither=1)
                quantized = quantized.convert('RGBA')
                quantized.save(output_buffer, format='PNG', optimize=True)
            else:
                quantized = input_image.convert('RGB').quantize(colors=colors, method=2, dither=1)
                quantized.save(output_buffer, format='PNG', optimize=True)
            output_format = 'PNG'
            mime = 'image/png'
        else:
            # Everything else: convert to WebP for excellent compression
            if input_image.mode == 'RGBA':
                input_image.save(output_buffer, format='WEBP', quality=quality, method=4)
            else:
                input_image.convert('RGB').save(output_buffer, format='WEBP', quality=quality, method=4)
            output_format = 'WEBP'
            mime = 'image/webp'

        compressed_bytes = output_buffer.getvalue()
        new_size = len(compressed_bytes)

        # Only use compressed if it's actually smaller
        if new_size >= original_size:
            img_b64 = base64.b64encode(raw_bytes).decode('utf-8')
            new_size = original_size
            mime = f'image/{original_format.lower()}'
        else:
            img_b64 = base64.b64encode(compressed_bytes).decode('utf-8')

        saved_pct = round((1 - new_size / original_size) * 100) if original_size > 0 else 0

        return jsonify({
            "data": img_b64,
            "mime": mime,
            "originalSize": original_size,
            "newSize": new_size,
            "savedPercentage": f"{saved_pct}%"
        })

    except Exception as e:
        print(f"Compression error: {e}")
        return jsonify({"error": str(e)}), 500

def _preprocess_binary(img):
    """
    Flatten transparency → RGB distance-from-white threshold.
    Uses Euclidean distance from pure white in RGB space so bright-coloured pixels
    (yellow, light orange, cyan, etc.) are correctly kept as foreground — unlike
    grayscale luminance which misclassifies them as near-white background.
    """
    # Composite transparent areas onto white background
    white = Image.new('RGBA', img.size, (255, 255, 255, 255))
    white.paste(img, mask=img.split()[3])
    rgb = white.convert('RGB')

    # Light blur to suppress JPEG/anti-aliasing noise before thresholding
    rgb = rgb.filter(ImageFilter.GaussianBlur(radius=0.8))

    arr = np.array(rgb, dtype=np.float32)
    # Distance from pure white in RGB space — yellow (255,255,0) → dist≈255, white → 0
    dist = np.sqrt(
        (255 - arr[:, :, 0]) ** 2 +
        (255 - arr[:, :, 1]) ** 2 +
        (255 - arr[:, :, 2]) ** 2
    )
    # Pixels with distance > threshold are foreground (black), rest are background (white)
    binary = np.where(dist > 50, 0, 255).astype(np.uint8)

    result = Image.fromarray(binary, 'L').convert('RGBA')
    return result

def _preprocess_color(img):
    """
    Flatten transparency → slight saturation boost so vtracer's color
    quantisation separates colours more cleanly.
    """
    white = Image.new('RGBA', img.size, (255, 255, 255, 255))
    white.paste(img, mask=img.split()[3])
    rgb = white.convert('RGB')
    rgb = ImageEnhance.Color(rgb).enhance(1.4)   # boost saturation 40%
    return rgb.convert('RGBA')

@app.route('/vectorize', methods=['POST'])
def vectorize():
    if 'image' not in request.files:
        return jsonify({"error": "No image file provided"}), 400

    file = request.files['image']
    colormode = request.form.get('colormode', 'color')  # 'color' or 'binary'
    mode = request.form.get('mode', 'spline')            # 'spline' | 'polygon' | 'none'
    hierarchical = request.form.get('hierarchical', 'stacked')  # 'stacked' | 'cutout'
    layer_difference = int(request.form.get('layer_difference', 16))
    corner_threshold = int(request.form.get('corner_threshold', 60))
    length_threshold = float(request.form.get('length_threshold', 4.0))
    splice_threshold = int(request.form.get('splice_threshold', 45))
    filter_speckle = int(request.form.get('filter_speckle', 4))
    color_precision = int(request.form.get('color_precision', 6))
    path_precision = int(request.form.get('path_precision', 8))

    # Validate enums so a bad value can't crash vtracer
    if mode not in ('spline', 'polygon', 'none'):
        mode = 'spline'
    if hierarchical not in ('stacked', 'cutout'):
        hierarchical = 'stacked'

    try:
        raw_bytes = file.read()

        # Normalize any image format (WebP, AVIF, JPEG, etc.) to PNG
        # so vtracer can always decode it
        img = Image.open(io.BytesIO(raw_bytes))
        if img.mode == 'RGBA':
            pass  # keep alpha
        else:
            img = img.convert('RGBA')

        # Cap large images for performance (2000px max side)
        img.thumbnail((2000, 2000), Image.LANCZOS)

        # Preprocess before vectorisation
        if colormode == 'binary':
            img = _preprocess_binary(img)
        else:
            img = _preprocess_color(img)

        png_buffer = io.BytesIO()
        img.save(png_buffer, format='PNG')
        png_bytes = png_buffer.getvalue()

        svg_str = vtracer.convert_raw_image_to_svg(
            png_bytes,
            img_format='png',
            colormode=colormode,
            hierarchical=hierarchical,
            mode=mode,
            filter_speckle=filter_speckle,
            color_precision=color_precision,
            layer_difference=layer_difference,
            corner_threshold=corner_threshold,
            length_threshold=length_threshold,
            splice_threshold=splice_threshold,
            path_precision=path_precision,
        )

        path_count = svg_str.count('<path')
        svg_size = len(svg_str.encode('utf-8'))

        return jsonify({
            "svg": svg_str,
            "colormode": colormode,
            "pathCount": path_count,
            "svgSize": svg_size,
        })

    except Exception as e:
        print(f"Vectorization error: {e}")
        return jsonify({"error": str(e)}), 500

# --- OCR Engine (RapidOCR via ONNX Runtime, lazy init) ---
_ocr_engine = None

def get_ocr_reader():
    global _ocr_engine
    if _ocr_engine is None:
        print("[OCR] Loading RapidOCR engine (first call)...")
        from rapidocr_onnxruntime import RapidOCR

        # Use Arabic recognition model if available (bundled in models/ocr/)
        _app_dir = os.path.dirname(os.path.abspath(__file__))
        _ar_rec = os.path.join(_app_dir, 'models', 'ocr', 'arabic_rec.onnx')
        _ar_dict = os.path.join(_app_dir, 'models', 'ocr', 'arabic_dict.txt')

        if os.path.isfile(_ar_rec) and os.path.isfile(_ar_dict):
            print(f"[OCR] Using Arabic rec model: {_ar_rec}")
            _ocr_engine = RapidOCR(rec_model_path=_ar_rec, rec_keys_path=_ar_dict)
        else:
            print("[OCR] Arabic model not found, using default (Chinese/English)")
            _ocr_engine = RapidOCR()

        print("[OCR] Engine ready.")
    return _ocr_engine

import re as _re

_ARABIC_RE = _re.compile(r'[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF\uFB50-\uFDFF\uFE70-\uFEFF]')
_NUMBERS_RE = _re.compile(r'[\d,.]+')

def _fix_arabic_line(text):
    """Fix reversed Arabic text from OCR (visual→logical order)."""
    if not _ARABIC_RE.search(text):
        return text  # No Arabic chars — leave as-is
    # Reverse the whole line (converts visual LTR → logical RTL)
    reversed_text = text[::-1]
    # Numbers got reversed too — flip them back
    reversed_text = _NUMBERS_RE.sub(lambda m: m.group(0)[::-1], reversed_text)
    return reversed_text

def _extract_text(result):
    """Extract plain text from RapidOCR result: [[bbox, text, confidence], ...]"""
    if not result:
        return ''
    return '\n'.join(_fix_arabic_line(line[1]) for line in result if line)

@app.route('/ocr', methods=['POST'])
def ocr():
    try:
        get_ocr_reader()
    except ImportError:
        return jsonify({"error": "OCR engine not available."}), 503

    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    filename = (file.filename or '').lower()
    raw_bytes = file.read()

    try:
        engine = get_ocr_reader()
        pages_text = []

        if filename.endswith('.pdf'):
            # Convert each PDF page to an image array, run OCR per page
            pdf_doc = fitz.open(stream=raw_bytes, filetype='pdf')
            for page_num in range(len(pdf_doc)):
                page = pdf_doc[page_num]
                mat = fitz.Matrix(2.0, 2.0)  # 2x scale for better OCR accuracy
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)
                result, _ = engine(img_array)
                page_text = _extract_text(result)
                pages_text.append(f"--- Page {page_num + 1} ---\n{page_text}")
            pdf_doc.close()
        else:
            # Image file
            img_array = np.array(Image.open(io.BytesIO(raw_bytes)).convert('RGB'))
            result, _ = engine(img_array)
            pages_text.append(_extract_text(result))

        full_text = '\n\n'.join(pages_text)
        return jsonify({
            "text": full_text,
            "pages": len(pages_text)
        })

    except Exception as e:
        print(f"[OCR] Error: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

# --- PDF Tools ---

@app.route('/pdf/thumbnails', methods=['POST'])
def pdf_thumbnails():
    if 'pdf' not in request.files:
        return jsonify({"error": "No PDF provided"}), 400

    file = request.files['pdf']
    dpi = int(request.form.get('dpi', 72))

    try:
        raw = file.read()
        doc = fitz.open(stream=raw, filetype='pdf')
        scale = dpi / 72
        mat = fitz.Matrix(scale, scale)

        thumbnails = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("png")
            b64 = base64.b64encode(img_bytes).decode('utf-8')
            thumbnails.append({
                "pageNum": page_num,
                "data": f"data:image/png;base64,{b64}",
                "width": pix.width,
                "height": pix.height
            })

        page_count = len(doc)
        doc.close()
        return jsonify({"thumbnails": thumbnails, "pageCount": page_count})

    except Exception as e:
        print(f"[PDF Thumbnails] Error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route('/preview/thumbnail', methods=['POST'])
def preview_thumbnail():
    """Generic file thumbnail — renders page 0 of PDF/AI via PyMuPDF,
    or opens PSD/TIFF/etc. via Pillow. Returns PNG as base64 data URL."""
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    max_size = int(request.form.get('maxSize', 128))
    filename = (file.filename or '').lower()
    raw = file.read()

    try:
        img = None

        # AI / PDF — render page 0 via PyMuPDF
        if filename.endswith('.ai') or filename.endswith('.pdf'):
            doc = fitz.open(stream=raw, filetype='pdf')
            if len(doc) > 0:
                page = doc[0]
                scale = min(max_size / page.rect.width, max_size / page.rect.height, 2.0)
                pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
                img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            doc.close()

        # Office OOXML — extract embedded thumbnail from ZIP archive
        elif filename.endswith(('.docx', '.xlsx', '.pptx')):
            try:
                zf = zipfile.ZipFile(io.BytesIO(raw))
                thumb_data = None
                for candidate in ['docProps/thumbnail.jpeg', 'docProps/thumbnail.png']:
                    if candidate in zf.namelist():
                        thumb_data = zf.read(candidate)
                        break
                zf.close()
                if thumb_data:
                    img = Image.open(io.BytesIO(thumb_data))
                    img.load()
            except zipfile.BadZipFile:
                pass

        else:
            # PSD, TIFF, and any other Pillow-supported format
            img = Image.open(io.BytesIO(raw))
            img.load()  # Force decode (important for PSD lazy loading)

        if img is None:
            return jsonify({"error": "Could not render file"}), 400

        # Composite alpha onto dark background (matches app UI)
        if img.mode in ('RGBA', 'LA', 'PA'):
            bg = Image.new('RGBA', img.size, (30, 30, 30, 255))
            bg.paste(img, mask=img.split()[-1])
            img = bg.convert('RGB')
        elif img.mode != 'RGB':
            img = img.convert('RGB')

        img.thumbnail((max_size, max_size), Image.LANCZOS)

        buf = io.BytesIO()
        img.save(buf, format='PNG', optimize=True)
        b64 = base64.b64encode(buf.getvalue()).decode('utf-8')

        return jsonify({
            "data": f"data:image/png;base64,{b64}",
            "width": img.width,
            "height": img.height,
        })

    except Exception as e:
        print(f"[Preview] Error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route('/pdf/merge', methods=['POST'])
def pdf_merge():
    files = request.files.getlist('pdfs')
    if len(files) < 2:
        return jsonify({"error": "Need at least 2 PDFs"}), 400

    try:
        merged = fitz.open()
        for f in files:
            raw = f.read()
            doc = fitz.open(stream=raw, filetype='pdf')
            merged.insert_pdf(doc)
            doc.close()

        page_count = len(merged)
        out_bytes = merged.tobytes(deflate=True, garbage=4)
        merged.close()

        b64 = base64.b64encode(out_bytes).decode('utf-8')
        return jsonify({
            "data": b64,
            "size": len(out_bytes),
            "pageCount": page_count
        })

    except Exception as e:
        print(f"[PDF Merge] Error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route('/pdf/organize', methods=['POST'])
def pdf_organize():
    import json
    files = request.files.getlist('pdfs')
    pages_json = request.form.get('pages', '[]')

    try:
        page_order = json.loads(pages_json)

        docs = []
        for f in files:
            raw = f.read()
            doc = fitz.open(stream=raw, filetype='pdf')
            docs.append(doc)

        result = fitz.open()
        for entry in page_order:
            fi = entry['fileIndex']
            pn = entry['pageNum']
            if fi < len(docs) and pn < len(docs[fi]):
                result.insert_pdf(docs[fi], from_page=pn, to_page=pn)

        page_count = len(result)
        out_bytes = result.tobytes(deflate=True, garbage=4)
        result.close()
        for d in docs:
            d.close()

        b64 = base64.b64encode(out_bytes).decode('utf-8')
        return jsonify({
            "data": b64,
            "size": len(out_bytes),
            "pageCount": page_count
        })

    except Exception as e:
        print(f"[PDF Organize] Error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route('/pdf/compress', methods=['POST'])
def pdf_compress():
    if 'pdf' not in request.files:
        return jsonify({"error": "No PDF provided"}), 400

    file = request.files['pdf']
    preset = request.form.get('preset', 'medium')

    presets = {
        'low':    {'image_quality': 30, 'dpi': 100},
        'medium': {'image_quality': 55, 'dpi': 150},
        'high':   {'image_quality': 80, 'dpi': 200},
    }
    params = presets.get(preset, presets['medium'])

    try:
        raw = file.read()
        original_size = len(raw)
        doc = fitz.open(stream=raw, filetype='pdf')

        # Re-compress images on each page
        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images(full=True)
            for img_info in image_list:
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    if not base_image or not base_image.get("image"):
                        continue
                    img_bytes = base_image["image"]
                    img = Image.open(io.BytesIO(img_bytes))

                    target_dpi = params['dpi']
                    scale = target_dpi / 150
                    if scale < 1:
                        new_w = max(1, int(img.width * scale))
                        new_h = max(1, int(img.height * scale))
                        img = img.resize((new_w, new_h), Image.LANCZOS)

                    buf = io.BytesIO()
                    if img.mode in ('RGBA', 'P'):
                        img = img.convert('RGB')
                    img.save(buf, format='JPEG', quality=params['image_quality'], optimize=True)
                    new_img_bytes = buf.getvalue()

                    if len(new_img_bytes) < len(img_bytes):
                        # Replace image in PDF
                        page.replace_image(xref, stream=new_img_bytes)
                except Exception:
                    continue

        out_bytes = doc.tobytes(deflate=True, garbage=4, clean=True)
        doc.close()

        new_size = len(out_bytes)
        if new_size >= original_size:
            b64 = base64.b64encode(raw).decode('utf-8')
            new_size = original_size
        else:
            b64 = base64.b64encode(out_bytes).decode('utf-8')

        saved_pct = round((1 - new_size / original_size) * 100) if original_size > 0 else 0

        return jsonify({
            "data": b64,
            "originalSize": original_size,
            "newSize": new_size,
            "savedPercentage": f"{saved_pct}%"
        })

    except Exception as e:
        print(f"[PDF Compress] Error: {e}")
        return jsonify({"error": str(e)}), 500


# Cache the Unicode font path used for the (invisible) OCR text layer.
_ocr_layer_font = None

def _get_ocr_layer_font():
    """Register a Unicode TTF (Arabic+Latin) for the searchable text layer.
    Returns the reportlab font name; falls back to 'Helvetica' (Latin only)."""
    global _ocr_layer_font
    if _ocr_layer_font is not None:
        return _ocr_layer_font
    _ocr_layer_font = 'Helvetica'
    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        candidates = [
            os.path.join(os.path.dirname(os.path.abspath(__file__)), 'models', 'ocr', 'ocr_layer.ttf'),
            r'C:\Windows\Fonts\arial.ttf',
            '/System/Library/Fonts/Supplemental/Arial Unicode.ttf',  # macOS 10.15+ (Arabic + Latin)
            '/Library/Fonts/Arial Unicode.ttf',                       # older macOS
            '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
        ]
        for fp in candidates:
            if os.path.isfile(fp):
                pdfmetrics.registerFont(TTFont('ocrlayer', fp))
                _ocr_layer_font = 'ocrlayer'
                break
    except Exception as e:
        print(f"[PDF Searchable] font register skipped: {e}")
    return _ocr_layer_font


@app.route('/pdf/searchable', methods=['POST'])
def pdf_searchable():
    """Add an invisible OCR text layer so a scanned PDF becomes selectable/searchable.
    Permissive stack: pypdfium2 (render) + RapidOCR (ocr) + reportlab (overlay) + pypdf (merge)."""
    if 'pdf' not in request.files:
        return jsonify({"error": "No PDF provided"}), 400
    try:
        get_ocr_reader()
    except ImportError:
        return jsonify({"error": "OCR engine not available."}), 503

    file = request.files['pdf']
    raw = file.read()

    try:
        import pypdfium2 as pdfium
        from reportlab.pdfgen import canvas as rl_canvas
        from pypdf import PdfReader, PdfWriter

        engine = get_ocr_reader()
        font_name = _get_ocr_layer_font()
        DPI = 200
        scale = DPI / 72.0

        # 1) Build an invisible text overlay (one page per source page).
        overlay_buf = io.BytesIO()
        c = rl_canvas.Canvas(overlay_buf)
        src = pdfium.PdfDocument(raw)
        n_pages = len(src)
        for i in range(n_pages):
            page = src[i]
            w_pt, h_pt = page.get_size()  # PDF points
            c.setPageSize((w_pt, h_pt))

            pil = page.render(scale=scale).to_pil().convert('RGB')
            result, _ = engine(np.array(pil))

            for line in (result or []):
                box, text = line[0], line[1]
                if not text or not text.strip():
                    continue
                text = _fix_arabic_line(text)
                xs = [p[0] for p in box]; ys = [p[1] for p in box]
                x_pt = min(xs) / scale
                box_w = max(1.0, (max(xs) - min(xs)) / scale)
                box_h = max(1.0, (max(ys) - min(ys)) / scale)
                baseline = h_pt - (max(ys) / scale)
                fs = box_h * 0.9
                to = c.beginText()
                to.setTextRenderMode(3)  # 3 = invisible
                to.setFont(font_name, fs)
                try:
                    sw = c.stringWidth(text, font_name, fs)
                    if sw > 0:
                        to.setHorizScale(100.0 * box_w / sw)  # stretch to align with glyphs
                except Exception:
                    pass
                to.setTextOrigin(x_pt, baseline)
                to.textLine(text)
                c.drawText(to)

            c.showPage()
        c.save()
        src.close()
        overlay_buf.seek(0)

        # 2) Stamp the text layer onto the original pages (keeps original visuals intact).
        reader = PdfReader(io.BytesIO(raw))
        overlay = PdfReader(overlay_buf)
        writer = PdfWriter()
        for i, pg in enumerate(reader.pages):
            if i < len(overlay.pages):
                try:
                    pg.merge_page(overlay.pages[i])
                except Exception:
                    pass
            writer.add_page(pg)

        out = io.BytesIO()
        writer.write(out)
        out_bytes = out.getvalue()
        b64 = base64.b64encode(out_bytes).decode('utf-8')
        return jsonify({"data": b64, "size": len(out_bytes), "pages": n_pages})

    except Exception as e:
        print(f"[PDF Searchable] Error: {e}")
        import traceback; traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route('/pdf/to-word', methods=['POST'])
def pdf_to_word():
    if 'pdf' not in request.files:
        return jsonify({"error": "No PDF provided"}), 400

    file = request.files['pdf']
    raw = file.read()
    tmp_dir = tempfile.mkdtemp()

    try:
        pdf_path = os.path.join(tmp_dir, 'input.pdf')
        docx_path = os.path.join(tmp_dir, 'output.docx')
        with open(pdf_path, 'wb') as f:
            f.write(raw)

        from pdf2docx import Converter
        cv = Converter(pdf_path)
        cv.convert(docx_path)
        cv.close()

        with open(docx_path, 'rb') as f:
            out_bytes = f.read()

        b64 = base64.b64encode(out_bytes).decode('utf-8')
        return jsonify({"data": b64, "size": len(out_bytes)})

    except Exception as e:
        print(f"[PDF to Word] Error: {e}")
        return jsonify({"error": str(e)}), 500
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


@app.route('/pdf/to-pptx', methods=['POST'])
def pdf_to_pptx():
    if 'pdf' not in request.files:
        return jsonify({"error": "No PDF provided"}), 400

    file = request.files['pdf']
    dpi = int(request.form.get('dpi', 200))
    raw = file.read()
    tmp_dir = tempfile.mkdtemp()

    try:
        doc = fitz.open(stream=raw, filetype='pdf')
        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation()
        scale = dpi / 72
        slide_count = len(doc)

        for page_num in range(slide_count):
            page = doc[page_num]
            mat = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            # PDF points → EMU (1 pt = 12700 EMU)
            page_w_emu = int(page.rect.width * 12700)
            page_h_emu = int(page.rect.height * 12700)
            prs.slide_width = Emu(page_w_emu)
            prs.slide_height = Emu(page_h_emu)

            img_path = os.path.join(tmp_dir, f'page_{page_num}.png')
            pix.save(img_path)

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            slide.shapes.add_picture(img_path, 0, 0, Emu(page_w_emu), Emu(page_h_emu))

        doc.close()
        pptx_path = os.path.join(tmp_dir, 'output.pptx')
        prs.save(pptx_path)

        with open(pptx_path, 'rb') as f:
            out_bytes = f.read()

        b64 = base64.b64encode(out_bytes).decode('utf-8')
        return jsonify({"data": b64, "size": len(out_bytes), "slideCount": slide_count})

    except Exception as e:
        print(f"[PDF to PPTX] Error: {e}")
        return jsonify({"error": str(e)}), 500
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# --- Format Converter ---

convert_jobs = {}
convert_jobs_lock = threading.Lock()

def _mac_brew_bin(name):
    """macOS: Finder/Dock-launched apps inherit launchd's minimal PATH (no /opt/homebrew,
    no /usr/local), so shutil.which misses Homebrew installs. Probe those locations."""
    if sys.platform == 'darwin':
        for p in (f'/opt/homebrew/bin/{name}', f'/usr/local/bin/{name}'):
            if os.path.isfile(p):
                return p
    return None

def find_ffmpeg():
    """Return path to ffmpeg binary — downloaded, bundled, or system."""
    exe = 'ffmpeg.exe' if sys.platform == 'win32' else 'ffmpeg'
    if TOOLS_DIR:
        downloaded = os.path.join(TOOLS_DIR, 'converter', exe)
        if os.path.isfile(downloaded):
            return downloaded
    bundled = os.path.join(os.path.dirname(__file__), 'bin', exe)
    if os.path.isfile(bundled):
        return bundled
    if shutil.which('ffmpeg'):
        return 'ffmpeg'
    return _mac_brew_bin('ffmpeg')

def find_ffprobe():
    """Return path to ffprobe binary — downloaded, bundled, or system."""
    exe = 'ffprobe.exe' if sys.platform == 'win32' else 'ffprobe'
    if TOOLS_DIR:
        downloaded = os.path.join(TOOLS_DIR, 'converter', exe)
        if os.path.isfile(downloaded):
            return downloaded
    bundled = os.path.join(os.path.dirname(__file__), 'bin', exe)
    if os.path.isfile(bundled):
        return bundled
    if shutil.which('ffprobe'):
        return 'ffprobe'
    return _mac_brew_bin('ffprobe')

def get_duration(ffprobe_path, filepath):
    """Get media duration in seconds using ffprobe."""
    try:
        result = subprocess.run(
            [ffprobe_path, '-v', 'quiet', '-show_entries', 'format=duration',
             '-of', 'default=noprint_wrappers=1:nokey=1', filepath],
            capture_output=True, text=True, timeout=10
        )
        return float(result.stdout.strip())
    except Exception:
        return 0

FFMPEG_PRESETS = {
    'mp4':  ['-c:v', 'libx264', '-preset', 'fast', '-crf', '23', '-c:a', 'aac'],
    'webm': ['-c:v', 'libvpx-vp9', '-crf', '30', '-b:v', '0', '-c:a', 'libopus'],
    'mov':  ['-c:v', 'libx264', '-c:a', 'aac'],
    'avi':  ['-c:v', 'libx264', '-c:a', 'mp3'],
    'mkv':  ['-c:v', 'libx264', '-c:a', 'aac'],
    'mp3':  ['-vn', '-c:a', 'libmp3lame', '-q:a', '2'],
    'wav':  ['-vn', '-c:a', 'pcm_s16le'],
    'ogg':  ['-vn', '-c:a', 'libvorbis', '-q:a', '5'],
    'gif':  ['-vf', 'fps=15,scale=480:-1:flags=lanczos', '-loop', '0'],
    # AV1 (modern, ~30% smaller than HEVC). Software libaom at cpu-used 8 + row-mt
    # keeps CPU encoding tolerable; muxed into an MP4 container.
    'av1':  ['-c:v', 'libaom-av1', '-crf', '30', '-b:v', '0', '-cpu-used', '8',
             '-row-mt', '1', '-pix_fmt', 'yuv420p', '-c:a', 'aac', '-movflags', '+faststart'],
}

# Targets that are a codec living in a different container → real output extension.
OUTPUT_CONTAINER = {'av1': 'mp4'}

def run_ffmpeg_job(job_id, ffmpeg_path, ffprobe_path, input_path, output_path, target_format):
    """Background thread: run FFmpeg and track progress."""
    try:
        duration = get_duration(ffprobe_path, input_path) if ffprobe_path else 0

        codec_args = FFMPEG_PRESETS.get(target_format, [])
        cmd = [ffmpeg_path, '-y', '-i', input_path, *codec_args, '-progress', 'pipe:1', output_path]

        stderr_chunks = []
        proc = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                                bufsize=0)

        # Drain stderr in a background thread to prevent deadlock
        def _drain_stderr():
            for chunk in iter(lambda: proc.stderr.read(4096), b''):
                stderr_chunks.append(chunk)
        stderr_thread = threading.Thread(target=_drain_stderr, daemon=True)
        stderr_thread.start()

        for raw_line in proc.stdout:
            line = raw_line.decode('utf-8', errors='replace').strip()
            # FFmpeg outputs both out_time_us and out_time_ms (both in microseconds)
            if duration > 0 and (line.startswith('out_time_us=') or line.startswith('out_time_ms=')):
                try:
                    us = int(line.split('=')[1])
                    pct = min(99, int((us / 1_000_000) / duration * 100))
                    with convert_jobs_lock:
                        if job_id in convert_jobs:
                            convert_jobs[job_id]['progress'] = pct
                except (ValueError, ZeroDivisionError):
                    pass

        proc.wait()
        stderr_thread.join(timeout=5)
        if proc.returncode != 0:
            stderr_text = b''.join(stderr_chunks).decode('utf-8', errors='replace')
            raise RuntimeError(f"FFmpeg failed (code {proc.returncode}): {stderr_text[:500]}")

        with open(output_path, 'rb') as f:
            out_bytes = f.read()

        mime_map = {
            'mp4': 'video/mp4', 'av1': 'video/mp4', 'webm': 'video/webm', 'mov': 'video/quicktime',
            'avi': 'video/x-msvideo', 'mkv': 'video/x-matroska',
            'mp3': 'audio/mpeg', 'wav': 'audio/wav', 'ogg': 'audio/ogg',
            'gif': 'image/gif',
        }
        mime = mime_map.get(target_format, 'application/octet-stream')
        b64 = base64.b64encode(out_bytes).decode('utf-8')
        data_url = f"data:{mime};base64,{b64}"

        with convert_jobs_lock:
            if job_id in convert_jobs:
                convert_jobs[job_id].update({
                    'status': 'done', 'progress': 100,
                    'dataUrl': data_url, 'size': len(out_bytes)
                })

    except Exception as e:
        print(f"[Convert] Job {job_id} error: {e}")
        with convert_jobs_lock:
            if job_id in convert_jobs:
                convert_jobs[job_id].update({'status': 'error', 'error': str(e)})

    finally:
        # Cleanup temp files
        try:
            if os.path.exists(input_path):
                os.remove(input_path)
            if os.path.exists(output_path):
                os.remove(output_path)
        except Exception:
            pass


@app.route('/convert/status', methods=['GET'])
def convert_status():
    return jsonify({'ffmpeg': find_ffmpeg() is not None})


@app.route('/convert/image', methods=['POST'])
def convert_image():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400

    file = request.files['file']
    target = request.form.get('format', 'png').lower()

    allowed = {'jpg', 'jpeg', 'png', 'webp', 'avif', 'bmp', 'tiff'}
    if target not in allowed:
        return jsonify({'error': f'Unsupported format: {target}'}), 400

    try:
        raw = file.read()
        filename = (file.filename or '').lower()

        # AI files are PDF-internally — render via PyMuPDF then hand to Pillow
        if filename.endswith('.ai'):
            doc = fitz.open(stream=raw, filetype='pdf')
            if len(doc) == 0:
                return jsonify({'error': 'AI file has no pages'}), 400
            page = doc[0]
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            doc.close()
        else:
            img = Image.open(io.BytesIO(raw))

        # Map target to Pillow format name
        fmt_map = {'jpg': 'JPEG', 'jpeg': 'JPEG', 'png': 'PNG', 'webp': 'WEBP', 'avif': 'AVIF', 'bmp': 'BMP', 'tiff': 'TIFF'}
        pil_format = fmt_map[target]

        # Handle mode conversion
        if pil_format == 'JPEG' and img.mode in ('RGBA', 'P', 'LA'):
            img = img.convert('RGB')
        elif pil_format in ('PNG', 'WEBP', 'AVIF') and img.mode not in ('RGB', 'RGBA'):
            img = img.convert('RGBA')
        elif pil_format in ('BMP', 'TIFF') and img.mode not in ('RGB', 'RGBA'):
            img = img.convert('RGB')

        buf = io.BytesIO()
        save_kwargs = {}
        if pil_format == 'JPEG':
            save_kwargs = {'quality': 90, 'optimize': True}
        elif pil_format == 'WEBP':
            save_kwargs = {'quality': 90, 'method': 4}
        elif pil_format == 'AVIF':
            save_kwargs = {'quality': 80, 'speed': 6}  # AV1 image: ~30-50% smaller than JPEG
        elif pil_format == 'PNG':
            save_kwargs = {'optimize': True}

        img.save(buf, format=pil_format, **save_kwargs)
        out_bytes = buf.getvalue()

        mime_map = {'jpg': 'image/jpeg', 'jpeg': 'image/jpeg', 'png': 'image/png',
                    'webp': 'image/webp', 'avif': 'image/avif', 'bmp': 'image/bmp', 'tiff': 'image/tiff'}
        mime = mime_map[target]
        b64 = base64.b64encode(out_bytes).decode('utf-8')

        return jsonify({
            'data': f'data:{mime};base64,{b64}',
            'format': target,
            'size': len(out_bytes)
        })

    except Exception as e:
        print(f"[Convert Image] Error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/convert/video', methods=['POST'])
def convert_video():
    ffmpeg_path = find_ffmpeg()
    if not ffmpeg_path:
        return jsonify({'error': 'FFmpeg not available'}), 503

    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400

    file = request.files['file']
    target = request.form.get('format', 'mp4').lower()

    allowed = set(FFMPEG_PRESETS.keys())
    if target not in allowed:
        return jsonify({'error': f'Unsupported format: {target}'}), 400

    try:
        job_id = str(uuid.uuid4())
        tmp_dir = tempfile.mkdtemp(prefix='convert_')

        # Save input file
        in_ext = os.path.splitext(file.filename or '')[1] or '.mp4'
        input_path = os.path.join(tmp_dir, f'input{in_ext}')
        file.save(input_path)

        output_path = os.path.join(tmp_dir, f'output.{OUTPUT_CONTAINER.get(target, target)}')

        with convert_jobs_lock:
            convert_jobs[job_id] = {
                'status': 'processing', 'progress': 0,
                'dataUrl': None, 'size': None, 'error': None
            }

        ffprobe_path = find_ffprobe()
        thread = threading.Thread(
            target=run_ffmpeg_job,
            args=(job_id, ffmpeg_path, ffprobe_path, input_path, output_path, target),
            daemon=True
        )
        thread.start()

        return jsonify({'jobId': job_id})

    except Exception as e:
        print(f"[Convert Video] Error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/convert/video/progress/<job_id>', methods=['GET'])
def convert_video_progress(job_id):
    with convert_jobs_lock:
        job = convert_jobs.get(job_id)

    if not job:
        return jsonify({'error': 'Job not found'}), 404

    result = {
        'status': job['status'],
        'progress': job['progress'],
    }

    if job['status'] == 'done':
        result['dataUrl'] = job['dataUrl']
        result['size'] = job['size']
        # Cleanup job from memory
        with convert_jobs_lock:
            convert_jobs.pop(job_id, None)
    elif job['status'] == 'error':
        result['error'] = job['error']
        with convert_jobs_lock:
            convert_jobs.pop(job_id, None)

    return jsonify(result)


# --- Image Upscaler (Real-ESRGAN) ---

upscale_jobs = {}
upscale_jobs_lock = threading.Lock()

def find_realesrgan():
    """Return path to realesrgan-ncnn-vulkan binary — downloaded, bundled, or system."""
    exe = 'realesrgan-ncnn-vulkan.exe' if sys.platform == 'win32' else 'realesrgan-ncnn-vulkan'
    if TOOLS_DIR:
        downloaded = os.path.join(TOOLS_DIR, 'upscaler', exe)
        if os.path.isfile(downloaded):
            return downloaded
    bundled = os.path.join(os.path.dirname(__file__), 'bin', exe)
    if os.path.isfile(bundled):
        return bundled
    if shutil.which('realesrgan-ncnn-vulkan'):
        return 'realesrgan-ncnn-vulkan'
    return _mac_brew_bin('realesrgan-ncnn-vulkan')

def find_models_dir():
    """Return path to the upscaler models directory."""
    if TOOLS_DIR:
        downloaded = os.path.join(TOOLS_DIR, 'upscaler', 'models')
        if os.path.isdir(downloaded):
            return downloaded
    bundled = os.path.join(os.path.dirname(__file__), 'bin', 'models')
    if os.path.isdir(bundled):
        return bundled
    return None


def run_upscale_job(job_id, exe_path, models_dir, input_path, output_path, scale, model_name):
    """Background thread: run Real-ESRGAN and track progress.

    Keeps the output file on disk (no base64) — the /upscale/result endpoint
    serves it as a binary download to avoid blowing up renderer memory.
    """
    try:
        cmd = [exe_path, '-i', input_path, '-o', output_path,
               '-s', str(scale), '-n', model_name, '-f', 'png']
        if models_dir:
            cmd += ['-m', models_dir]

        stderr_chunks = []
        proc = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, bufsize=0)

        # Drain stderr in a background thread (progress info comes on stderr)
        def _drain_stderr():
            for chunk in iter(lambda: proc.stderr.read(256), b''):
                stderr_chunks.append(chunk)
                # Parse progress from stderr: "xx.xx%"
                text = chunk.decode('utf-8', errors='replace')
                import re
                pcts = re.findall(r'(\d+(?:\.\d+)?)%', text)
                if pcts:
                    try:
                        pct = min(99, int(float(pcts[-1])))
                        with upscale_jobs_lock:
                            if job_id in upscale_jobs:
                                upscale_jobs[job_id]['progress'] = pct
                    except (ValueError, IndexError):
                        pass
        stderr_thread = threading.Thread(target=_drain_stderr, daemon=True)
        stderr_thread.start()

        # stdout doesn't produce much — just drain it
        proc.stdout.read()
        proc.wait()
        stderr_thread.join(timeout=10)

        if proc.returncode != 0:
            stderr_text = b''.join(stderr_chunks).decode('utf-8', errors='replace')
            raise RuntimeError(f"Real-ESRGAN failed (code {proc.returncode}): {stderr_text[:500]}")

        out_size = os.path.getsize(output_path)

        with upscale_jobs_lock:
            if job_id in upscale_jobs:
                upscale_jobs[job_id].update({
                    'status': 'done', 'progress': 100,
                    'outputPath': output_path, 'size': out_size
                })

        # Clean up input only — output is served by /upscale/result then cleaned
        try:
            if os.path.exists(input_path):
                os.remove(input_path)
        except Exception:
            pass

    except Exception as e:
        print(f"[Upscale] Job {job_id} error: {e}")
        with upscale_jobs_lock:
            if job_id in upscale_jobs:
                upscale_jobs[job_id].update({'status': 'error', 'error': str(e)})
        # Clean everything on error
        try:
            if os.path.exists(input_path):
                os.remove(input_path)
            if os.path.exists(output_path):
                os.remove(output_path)
            tmp_dir = os.path.dirname(input_path)
            if os.path.isdir(tmp_dir) and tmp_dir.startswith(tempfile.gettempdir()):
                shutil.rmtree(tmp_dir, ignore_errors=True)
        except Exception:
            pass


@app.route('/upscale/status', methods=['GET'])
def upscale_status():
    return jsonify({'available': find_realesrgan() is not None})


@app.route('/upscale', methods=['POST'])
def upscale_image():
    exe_path = find_realesrgan()
    if not exe_path:
        return jsonify({'error': 'Real-ESRGAN not available'}), 503

    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400

    file = request.files['file']
    scale = int(request.form.get('scale', 4))
    model = request.form.get('model', 'realesrgan-x4plus')

    if scale not in (2, 3, 4):
        return jsonify({'error': f'Invalid scale: {scale}'}), 400

    allowed_models = {'realesrgan-x4plus', 'realesrgan-x4plus-anime'}
    if model not in allowed_models:
        return jsonify({'error': f'Invalid model: {model}'}), 400

    try:
        job_id = str(uuid.uuid4())
        tmp_dir = tempfile.mkdtemp(prefix='upscale_')

        in_ext = os.path.splitext(file.filename or '')[1] or '.png'
        input_path = os.path.join(tmp_dir, f'input{in_ext}')
        file.save(input_path)

        output_path = os.path.join(tmp_dir, 'output.png')
        models_dir = find_models_dir()

        with upscale_jobs_lock:
            upscale_jobs[job_id] = {
                'status': 'processing', 'progress': 0,
                'outputPath': None, 'size': None, 'error': None
            }

        thread = threading.Thread(
            target=run_upscale_job,
            args=(job_id, exe_path, models_dir, input_path, output_path, scale, model),
            daemon=True
        )
        thread.start()

        return jsonify({'jobId': job_id})

    except Exception as e:
        print(f"[Upscale] Error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/upscale/progress/<job_id>', methods=['GET'])
def upscale_progress(job_id):
    with upscale_jobs_lock:
        job = upscale_jobs.get(job_id)

    if not job:
        return jsonify({'error': 'Job not found'}), 404

    result = {
        'status': job['status'],
        'progress': job['progress'],
    }

    if job['status'] == 'done':
        result['size'] = job['size']
        # Don't remove job — /upscale/result/<job_id> still needs to serve the file
    elif job['status'] == 'error':
        result['error'] = job['error']
        with upscale_jobs_lock:
            upscale_jobs.pop(job_id, None)

    return jsonify(result)


@app.route('/upscale/result/<job_id>', methods=['GET'])
def upscale_result(job_id):
    """Serve the upscaled image as a binary file.

    Does NOT auto-cleanup — the frontend calls /upscale/cleanup/<job_id>
    when the user removes the item, clears the list, or closes the tool.
    This avoids holding any image data in the Electron renderer's memory.
    """
    with upscale_jobs_lock:
        job = upscale_jobs.get(job_id)

    if not job:
        return jsonify({'error': 'Job not found'}), 404

    if job['status'] != 'done' or not job.get('outputPath'):
        return jsonify({'error': 'Result not ready'}), 409

    output_path = job['outputPath']
    if not os.path.isfile(output_path):
        with upscale_jobs_lock:
            upscale_jobs.pop(job_id, None)
        return jsonify({'error': 'Result file missing'}), 410

    return send_file(output_path, mimetype='image/png')


@app.route('/upscale/cleanup/<job_id>', methods=['POST'])
def upscale_cleanup(job_id):
    """Explicitly clean up a finished upscale job's temp files."""
    with upscale_jobs_lock:
        job = upscale_jobs.pop(job_id, None)

    if not job:
        return jsonify({'ok': True})   # already gone — fine

    output_path = job.get('outputPath')
    if output_path:
        try:
            if os.path.exists(output_path):
                os.remove(output_path)
            tmp_dir = os.path.dirname(output_path)
            if os.path.isdir(tmp_dir) and tmp_dir.startswith(tempfile.gettempdir()):
                shutil.rmtree(tmp_dir, ignore_errors=True)
        except Exception:
            pass

    return jsonify({'ok': True})


# ── Metadata Scrubber ───────────────────────────────────────────────

@app.route('/scrub-metadata', methods=['POST'])
def scrub_metadata():
    """Strip EXIF / document metadata from images and PDFs."""
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    filename = (file.filename or '').lower()
    raw = file.read()
    original_size = len(raw)

    try:
        removed_fields = {}

        if filename.endswith('.pdf'):
            doc = fitz.open(stream=raw, filetype='pdf')
            old_meta = doc.metadata or {}
            removed_fields = {k: v for k, v in old_meta.items() if v}
            doc.set_metadata({})
            try:
                doc.del_xml_metadata()
            except Exception:
                pass
            out_bytes = doc.tobytes(deflate=True, garbage=4)
            doc.close()
            mime = 'application/pdf'
        else:
            img = Image.open(io.BytesIO(raw))
            fmt = (img.format or 'PNG').upper()

            # Collect EXIF info before removal
            from PIL.ExifTags import TAGS
            exif = img.getexif()
            if exif:
                for tag_id, value in exif.items():
                    tag_name = TAGS.get(tag_id, str(tag_id))
                    try:
                        removed_fields[tag_name] = str(value)[:100]
                    except Exception:
                        removed_fields[tag_name] = '<binary>'

            # Also check img.info for other metadata
            for key in ('exif', 'icc_profile', 'xmp', 'photoshop'):
                if key in img.info:
                    if key not in removed_fields:
                        removed_fields[key] = 'present'

            # Strip all metadata: create new image from pixel data only
            clean = Image.new(img.mode, img.size)
            clean.putdata(list(img.getdata()))

            buf = io.BytesIO()
            save_kwargs = {'format': fmt, 'optimize': True}
            if fmt == 'JPEG':
                save_kwargs['quality'] = 95
            elif fmt == 'PNG':
                save_kwargs['optimize'] = True
            clean.save(buf, **save_kwargs)
            out_bytes = buf.getvalue()

            fmt_lower = fmt.lower()
            if fmt_lower in ('jpg', 'jpeg'):
                mime = 'image/jpeg'
            else:
                mime = f'image/{fmt_lower}'

        b64 = base64.b64encode(out_bytes).decode('utf-8')
        return jsonify({
            "data": b64,
            "mime": mime,
            "originalSize": original_size,
            "newSize": len(out_bytes),
            "removedFields": removed_fields,
        })

    except Exception as e:
        print(f"[Metadata] Error: {e}")
        return jsonify({"error": str(e)}), 500


# ── Color Palette ───────────────────────────────────────────────────

@app.route('/extract-palette', methods=['POST'])
def extract_palette():
    """Extract dominant colors from an image."""
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    count = int(request.form.get('count', 8))
    count = max(2, min(count, 16))

    try:
        raw = file.read()
        img = Image.open(io.BytesIO(raw)).convert('RGB')

        # Resize for speed (max 150px on longest side)
        img.thumbnail((150, 150), Image.LANCZOS)

        # Quantize to more colors than needed, then filter by distance
        oversample = min(count * 3, 32)
        quantized = img.quantize(colors=oversample, method=Image.Quantize.MEDIANCUT)
        palette_data = quantized.getpalette()  # flat [R,G,B,R,G,B,...]
        pixels = list(quantized.getdata())
        total = len(pixels)

        # Count frequency of each color index
        freq = {}
        for idx in pixels:
            freq[idx] = freq.get(idx, 0) + 1

        # Build candidate list sorted by frequency (most dominant first)
        candidates = []
        for idx, count_val in sorted(freq.items(), key=lambda x: -x[1]):
            if idx * 3 + 2 >= len(palette_data):
                continue
            r = palette_data[idx * 3]
            g = palette_data[idx * 3 + 1]
            b = palette_data[idx * 3 + 2]
            percentage = round((count_val / total) * 100, 1)
            candidates.append({"rgb": (r, g, b), "percentage": percentage})

        # Greedy selection: pick most dominant, then skip colors too close
        MIN_DIST = 35  # minimum Euclidean distance in RGB space
        selected = []
        for c in candidates:
            if len(selected) >= count:
                break
            too_close = False
            for s in selected:
                dr = c["rgb"][0] - s["rgb"][0]
                dg = c["rgb"][1] - s["rgb"][1]
                db = c["rgb"][2] - s["rgb"][2]
                if (dr*dr + dg*dg + db*db) ** 0.5 < MIN_DIST:
                    too_close = True
                    break
            if not too_close:
                selected.append(c)

        colors = []
        for c in selected:
            r, g, b = c["rgb"]
            colors.append({
                "hex": f'#{r:02X}{g:02X}{b:02X}',
                "rgb": [r, g, b],
                "percentage": c["percentage"],
            })

        return jsonify({"colors": colors})

    except Exception as e:
        print(f"[Palette] Error: {e}")
        return jsonify({"error": str(e)}), 500


# ── Watermark ───────────────────────────────────────────────────────

@app.route('/watermark', methods=['POST'])
def watermark():
    """Add text watermark to images and PDFs."""
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    text = request.form.get('text', 'Watermark')
    opacity = int(request.form.get('opacity', 30))
    font_size = int(request.form.get('fontSize', 36))
    style = request.form.get('style', 'diagonal')  # diagonal | center | corner
    color_hex = request.form.get('color', '#888888')
    filename = (file.filename or '').lower()
    raw = file.read()

    # Parse hex color
    try:
        color_hex = color_hex.lstrip('#')
        cr = int(color_hex[0:2], 16)
        cg = int(color_hex[2:4], 16)
        cb = int(color_hex[4:6], 16)
    except Exception:
        cr, cg, cb = 136, 136, 136

    try:
        if filename.endswith('.pdf'):
            doc = fitz.open(stream=raw, filetype='pdf')
            color_fitz = (cr / 255, cg / 255, cb / 255)

            for page in doc:
                rect = page.rect
                if style == 'diagonal':
                    step_x = font_size * max(len(text), 4) * 0.6 + 80
                    step_y = font_size * 3
                    y = 0
                    while y < rect.height + rect.width:
                        x = -rect.width * 0.5
                        while x < rect.width * 1.5:
                            page.insert_text(
                                (x, y), text,
                                fontsize=font_size,
                                fontname="helv",
                                color=color_fitz,
                                rotate=45,
                                overlay=True,
                                fill_opacity=opacity / 100,
                            )
                            x += step_x
                        y += step_y
                elif style == 'center':
                    tw = fitz.get_text_length(text, fontname="helv", fontsize=font_size * 2)
                    x = (rect.width - tw) / 2
                    y = rect.height / 2
                    page.insert_text(
                        (x, y), text,
                        fontsize=font_size * 2,
                        fontname="helv",
                        color=color_fitz,
                        overlay=True,
                        fill_opacity=opacity / 100,
                    )
                else:  # corner
                    tw = fitz.get_text_length(text, fontname="helv", fontsize=font_size)
                    x = rect.width - tw - 20
                    y = rect.height - 20
                    page.insert_text(
                        (x, y), text,
                        fontsize=font_size,
                        fontname="helv",
                        color=color_fitz,
                        overlay=True,
                        fill_opacity=opacity / 100,
                    )

            out_bytes = doc.tobytes(deflate=True)
            doc.close()
            b64 = base64.b64encode(out_bytes).decode('utf-8')
            return jsonify({"data": b64, "mime": "application/pdf", "size": len(out_bytes)})

        else:
            # Image watermark via Pillow
            from PIL import ImageDraw, ImageFont
            import math

            img = Image.open(io.BytesIO(raw)).convert('RGBA')
            w, h = img.size

            # Load font — try system fonts with Arabic support
            font = None
            for font_path in [
                'C:/Windows/Fonts/tahoma.ttf',
                'C:/Windows/Fonts/arial.ttf',
                '/System/Library/Fonts/Supplemental/Arial Unicode.ttf',  # macOS (Arabic + Latin)
                '/System/Library/Fonts/Supplemental/Arial.ttf',          # macOS
                '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
            ]:
                try:
                    font = ImageFont.truetype(font_path, font_size)
                    break
                except Exception:
                    continue
            if font is None:
                try:
                    font = ImageFont.load_default(size=font_size)  # Pillow >= 10.1 honors size
                except TypeError:
                    font = ImageFont.load_default()

            alpha = int(255 * opacity / 100)
            color = (cr, cg, cb, alpha)

            watermark_layer = Image.new('RGBA', (w, h), (0, 0, 0, 0))
            draw = ImageDraw.Draw(watermark_layer)

            if style == 'diagonal':
                # Measure text
                bbox = draw.textbbox((0, 0), text, font=font)
                tw = bbox[2] - bbox[0]
                th = bbox[3] - bbox[1]

                # Create a single text tile, then rotate
                tile_w = tw + 80
                tile_h = th + 20
                tile = Image.new('RGBA', (tile_w, tile_h), (0, 0, 0, 0))
                tile_draw = ImageDraw.Draw(tile)
                tile_draw.text((0, 0), text, fill=color, font=font)
                rotated = tile.rotate(35, expand=True, resample=Image.BICUBIC)
                rw, rh = rotated.size

                # Tile across the image
                for ty in range(-rh, h + rh, rh + 30):
                    for tx in range(-rw, w + rw, rw + 40):
                        watermark_layer.paste(rotated, (tx, ty), rotated)

            elif style == 'center':
                big_font = None
                for font_path in [
                    'C:/Windows/Fonts/tahoma.ttf',
                    'C:/Windows/Fonts/arial.ttf',
                    '/System/Library/Fonts/Supplemental/Arial Unicode.ttf',  # macOS
                    '/System/Library/Fonts/Supplemental/Arial.ttf',
                    '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
                ]:
                    try:
                        big_font = ImageFont.truetype(font_path, font_size * 2)
                        break
                    except Exception:
                        continue
                if big_font is None:
                    big_font = font

                bbox = draw.textbbox((0, 0), text, font=big_font)
                tw = bbox[2] - bbox[0]
                th = bbox[3] - bbox[1]
                x = (w - tw) // 2
                y = (h - th) // 2
                draw.text((x, y), text, fill=color, font=big_font)

            else:  # corner (bottom-right)
                bbox = draw.textbbox((0, 0), text, font=font)
                tw = bbox[2] - bbox[0]
                th = bbox[3] - bbox[1]
                x = w - tw - 20
                y = h - th - 20
                draw.text((x, y), text, fill=color, font=font)

            result = Image.alpha_composite(img, watermark_layer)

            # Save as PNG to preserve quality
            buf = io.BytesIO()
            result.save(buf, format='PNG')
            out_bytes = buf.getvalue()

            b64 = base64.b64encode(out_bytes).decode('utf-8')
            return jsonify({"data": b64, "mime": "image/png", "size": len(out_bytes)})

    except Exception as e:
        print(f"[Watermark] Error: {e}")
        return jsonify({"error": str(e)}), 500


# ── Resize Image (Pillow) ────────────────────────────────────────────
@app.route('/resize', methods=['POST'])
def resize_image():
    if 'image' not in request.files:
        return jsonify({"error": "No image"}), 400
    file = request.files['image']
    width = request.form.get('width', type=int)
    height = request.form.get('height', type=int)
    mode = request.form.get('mode', 'fit')   # 'fit' | 'exact' | 'width'
    max_kb = request.form.get('maxKB', type=int)
    try:
        img = Image.open(io.BytesIO(file.read()))
        fmt = (img.format or 'PNG').upper()
        ow, oh = img.size

        if mode == 'fill' and width and height:
            # Aspect ratio: cover-resize then centre-crop. NEVER enlarge — if the requested
            # box exceeds what the source can supply at this ratio, shrink it to fit.
            tw, th = width, height
            fit = min(ow / tw, oh / th, 1.0)
            if fit < 1.0:
                tw, th = max(1, round(tw * fit)), max(1, round(th * fit))
            scale = max(tw / ow, th / oh)
            rw, rh = max(1, round(ow * scale)), max(1, round(oh * scale))
            img = img.resize((rw, rh), Image.LANCZOS)
            left, top = (rw - tw) // 2, (rh - th) // 2
            img = img.crop((left, top, left + tw, top + th))
        elif mode == 'exact' and width and height:
            tw, th = width, height
            if (tw, th) != (ow, oh):
                img = img.resize((tw, th), Image.LANCZOS)
        elif mode == 'width' and width:
            tw = min(width, ow)  # never enlarge — cap at source width (use the Upscaler to go bigger)
            th = max(1, round(oh * (tw / ow)))
            if (tw, th) != (ow, oh):
                img = img.resize((tw, th), Image.LANCZOS)
        else:  # fit within width/height, keep aspect
            bw, bh = width or ow, height or oh
            ratio = min(bw / ow, bh / oh)
            tw, th = max(1, round(ow * ratio)), max(1, round(oh * ratio))
            if (tw, th) != (ow, oh):
                img = img.resize((tw, th), Image.LANCZOS)

        save_fmt = fmt if fmt in ('JPEG', 'PNG', 'WEBP') else 'PNG'
        if save_fmt == 'JPEG' and img.mode in ('RGBA', 'P', 'LA'):
            img = img.convert('RGB')

        out = io.BytesIO()
        if max_kb and save_fmt in ('JPEG', 'WEBP'):
            q = 90
            while q >= 20:
                out.seek(0); out.truncate()
                img.save(out, format=save_fmt, quality=q)
                if out.tell() <= max_kb * 1024:
                    break
                q -= 8
        elif save_fmt in ('JPEG', 'WEBP'):
            img.save(out, format=save_fmt, quality=90)
        else:
            img.save(out, format=save_fmt, optimize=True)

        data = out.getvalue()
        mime = {'JPEG': 'image/jpeg', 'PNG': 'image/png', 'WEBP': 'image/webp'}[save_fmt]
        return jsonify({"data": base64.b64encode(data).decode('utf-8'),
                        "mime": mime, "width": tw, "height": th, "size": len(data)})
    except Exception as e:
        print(f"[Resize] Error: {e}")
        return jsonify({"error": str(e)}), 500


# ── Zip / Unzip (stdlib zipfile) ──────────────────────────────────────
@app.route('/zip', methods=['POST'])
def zip_files():
    files = request.files.getlist('files')
    if not files:
        return jsonify({"error": "No files"}), 400
    try:
        buf = io.BytesIO()
        seen = {}
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as z:
            for f in files:
                name = os.path.basename(f.filename or 'file')
                if name in seen:
                    seen[name] += 1
                    base, ext = os.path.splitext(name)
                    name = f"{base} ({seen[name]}){ext}"
                else:
                    seen[name] = 0
                z.writestr(name, f.read())
        data = buf.getvalue()
        return jsonify({"data": base64.b64encode(data).decode('utf-8'), "size": len(data)})
    except Exception as e:
        print(f"[Zip] Error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route('/unzip', methods=['POST'])
def unzip_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file"}), 400
    f = request.files['file']
    try:
        results = []
        with zipfile.ZipFile(io.BytesIO(f.read())) as z:
            for info in z.infolist():
                if info.is_dir():
                    continue
                with z.open(info) as zf:
                    data = zf.read()
                results.append({"name": os.path.basename(info.filename) or info.filename,
                                "data": base64.b64encode(data).decode('utf-8')})
        return jsonify({"results": results, "count": len(results)})
    except zipfile.BadZipFile:
        return jsonify({"error": "Not a valid zip file"}), 400
    except Exception as e:
        print(f"[Unzip] Error: {e}")
        return jsonify({"error": str(e)}), 500


# ── PDF ↔ Images (pypdfium2 render / Pillow combine) ──────────────────
@app.route('/pdf/to-images', methods=['POST'])
def pdf_to_images():
    if 'pdf' not in request.files:
        return jsonify({"error": "No PDF"}), 400
    f = request.files['pdf']
    dpi = request.form.get('dpi', 150, type=int)
    try:
        import pypdfium2 as pdfium
        base = os.path.splitext(os.path.basename(f.filename or 'page'))[0]
        doc = pdfium.PdfDocument(f.read())
        scale = dpi / 72.0
        results = []
        for i in range(len(doc)):
            pil = doc[i].render(scale=scale).to_pil().convert('RGB')
            out = io.BytesIO(); pil.save(out, format='PNG')
            results.append({"name": f"{base}-{i + 1}.png",
                            "data": base64.b64encode(out.getvalue()).decode('utf-8')})
        doc.close()
        return jsonify({"results": results, "count": len(results)})
    except Exception as e:
        print(f"[PDF->Images] Error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route('/pdf/from-images', methods=['POST'])
def pdf_from_images():
    files = request.files.getlist('images')
    if not files:
        return jsonify({"error": "No images"}), 400
    try:
        imgs = []
        for f in files:
            im = Image.open(io.BytesIO(f.read()))
            imgs.append(im.convert('RGB') if im.mode != 'RGB' else im)
        out = io.BytesIO()
        imgs[0].save(out, format='PDF', save_all=True, append_images=imgs[1:])
        data = out.getvalue()
        return jsonify({"data": base64.b64encode(data).decode('utf-8'),
                        "size": len(data), "pages": len(imgs)})
    except Exception as e:
        print(f"[Images->PDF] Error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route('/tools/status', methods=['GET'])
def tools_status():
    """Report which tool backends have their deps satisfied."""
    status = {}

    # Remover (BEN2): check if the ben2 package is present (without importing torch)
    import importlib.util
    status['remover'] = importlib.util.find_spec('ben2') is not None

    # OCR: check if rapidocr is importable
    try:
        import rapidocr_onnxruntime
        status['ocr'] = True
    except ImportError:
        status['ocr'] = False

    # Upscaler: check if binary exists
    status['upscaler'] = find_realesrgan() is not None

    # Converter: check if ffmpeg exists
    status['converter'] = find_ffmpeg() is not None

    # Default tools (Pillow, PyMuPDF, vtracer ship with app)
    for tid in ['compressor', 'cropper', 'vectorizer', 'pdf', 'metadata', 'watermark', 'palette', 'resize', 'zip']:
        status[tid] = True

    # Shelf is pure Electron
    status['shelf'] = True

    return jsonify(status)


def _preload_models():
    """Background preload of light, bundled engines so first use feels instant.
    (The remover/BEN2 is heavy + on-demand, so it is NOT preloaded — it loads on
    first drop and unloads after idle.)"""
    try:
        get_ocr_reader()
        print("[Preload] OCR engine ready.")
    except Exception as e:
        print(f"[Preload] OCR skip: {e}")


if __name__ == '__main__':
    # Preload models in background thread so first use is fast
    threading.Thread(target=_preload_models, daemon=True).start()

    # Dev-only memory logger
    if os.environ.get('DRAGIN_DEV_MEMORY_LOG'):
        import psutil
        def _mem_log():
            p = psutil.Process()
            while True:
                m = p.memory_info()
                print(f"[Memory:Python] RSS: {m.rss // 1024 // 1024}MB | VMS: {m.vms // 1024 // 1024}MB", flush=True)
                threading.Event().wait(30)
        threading.Thread(target=_mem_log, daemon=True).start()
    # Port 5000 is taken by macOS Control Center's AirPlay Receiver by default, so use an
    # uncommon fixed port instead (must match BASE_URL in src/services/api.ts).
    app.run(host='127.0.0.1', port=8756)
